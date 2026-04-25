"""
summarize_lecture.py
====================
Riassume una lezione universitaria combinando:
  - il testo (e le immagini) estratto da un .pptx
  - la trascrizione del professore (uno o più file .txt con timestamp float)

Dipendenze:
    pip install python-pptx ollama Pillow

Uso:
    # Riassunto completo (slide + trascrizione)
    python summarize_lecture.py slides.pptx trascrizione.txt
    python summarize_lecture.py slides.pptx parte1.txt parte2.txt parte3.txt

    # Solo riassunto delle slide (nessun file di trascrizione)
    python summarize_lecture.py slides.pptx --slides-only

    # Opzioni
    python summarize_lecture.py slides.pptx trascrizione.txt \\
        --model gemma4:e4b \\
        --output riassunto.md \\
        --no-images \\
        --verbose

Formato trascrizione atteso (.txt):
    [0.0] Benvenuti alla lezione di oggi.
    [4.85] Oggi parleremo di reti neurali convoluzionali.
    [12.3] Come vedete in questa figura...
"""

import argparse
import base64
import io
import re
import sys
from pathlib import Path

import ollama
import pdfplumber
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

PIL_AVAILABLE = True


# ─────────────────────────────────────────────
#  Estrazione testo dal .pptx
# ─────────────────────────────────────────────

def extract_slide_text(pptx_path: Path) -> str:
    """
    Estrae il testo di tutte le slide in un unico blocco markdown.
    Produce sezioni "### Slide N: Titolo".
    """
    prs = Presentation(str(pptx_path))
    sections = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        bodies = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if (hasattr(shape, "placeholder_format")
                    and shape.placeholder_format
                    and shape.placeholder_format.idx in (0, 1)):
                title = text
            else:
                bodies.append(text)

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        header = f"### Slide {idx}" + (f": {title}" if title else "")
        parts = [header]
        if bodies:
            parts.append("\n".join(bodies))
        if notes:
            parts.append(f"[Note: {notes}]")

        sections.append("\n".join(parts))

    return "\n\n".join(sections)


# ─────────────────────────────────────────────
#  Estrazione immagini dal .pptx
# ─────────────────────────────────────────────

def extract_slide_images(pptx_path: Path, max_images: int = 20) -> list[dict]:
    """
    Estrae le immagini dalle slide come JPEG base64.
    Salta i duplicati (stessa immagine in più slide).
    Limita a max_images totali per non sovraccaricare il contesto.
    """
    if not PIL_AVAILABLE:
        print("⚠️  Pillow non installato: immagini saltate (pip install Pillow)")
        return []

    prs = Presentation(str(pptx_path))
    images = []
    seen_hashes: set[int] = set()

    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type != 13:  # 13 = PICTURE
                continue
            try:
                blob = shape.image.blob
                blob_hash = hash(blob[:512])
                if blob_hash in seen_hashes:
                    continue
                seen_hashes.add(blob_hash)

                img = Image.open(io.BytesIO(blob))
                img.thumbnail((800, 800), Image.Resampling.LANCZOS)
                if img.mode not in ("RGB", "L"):
                    img = img.convert("RGB")

                buf = io.BytesIO()
                img.save(buf, format="JPEG", quality=75)
                b64 = base64.b64encode(buf.getvalue()).decode()

                images.append({"slide_idx": idx, "b64": b64})

                if len(images) >= max_images:
                    print(f"ℹ️  Limite {max_images} immagini raggiunto, le successive sono saltate.")
                    return images
            except Exception as e:
                print(f"  ⚠️  Slide {idx}: impossibile estrarre immagine ({e})")

    return images

# ─────────────────────────────────────────────
#  Estrazione testo e immagini dal .pdf
# ─────────────────────────────────────────────

def extract_slide_text_pdf(pdf_path: Path) -> str:
    """
    Estrae il testo pagina per pagina da un PDF.
    Usa pdfplumber che preserva meglio l'ordine di lettura.
    Produce sezioni "### Pagina N" analoghe a "### Slide N" del pptx.
    """

    sections = []
    with pdfplumber.open(str(pdf_path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            text = text.strip()
            if not text:
                text = "(pagina senza testo estraibile)"
            sections.append(f"### Pagina {idx}\n{text}")

    return "\n\n".join(sections)


def extract_slide_images_pdf(pdf_path: Path, max_images: int = 20) -> list[dict]:
    """
    Estrae le immagini embedded da un PDF come JPEG base64.
    Salta i duplicati tramite hash parziale del blob.
    Limita a max_images totali.
    Interfaccia identica a extract_slide_images().
    """
    if not PIL_AVAILABLE:
        print("⚠️  Pillow non installato: immagini saltate (pip install Pillow)")
        return []

    reader = PdfReader(str(pdf_path))
    images = []
    seen_hashes: set[int] = set()

    for idx, page in enumerate(reader.pages, start=1):
        # pypdf >= 3.x espone page.images come lista di ImageFile
        for img_obj in page.images:
            try:
                blob = img_obj.data
                blob_hash = hash(blob[:512])
                if blob_hash in seen_hashes:
                    continue
                seen_hashes.add(blob_hash)

                img = Image.open(io.BytesIO(blob))
                img.thumbnail((800, 800), Image.Resampling.LANCZOS)
                if img.mode not in ("RGB", "L"):
                    img = img.convert("RGB")

                buf = io.BytesIO()
                img.save(buf, format="JPEG", quality=75)
                b64 = base64.b64encode(buf.getvalue()).decode()

                images.append({"slide_idx": idx, "b64": b64})

                if len(images) >= max_images:
                    print(f"ℹ️  Limite {max_images} immagini raggiunto, le successive sono saltate.")
                    return images
            except Exception as e:
                print(f"  ⚠️  Pagina {idx}: impossibile estrarre immagine ({e})")

    return images

# ─────────────────────────────────────────────
#  Parsing e merge delle trascrizioni
# ─────────────────────────────────────────────

TRANSCRIPT_LINE = re.compile(r"\[(\d+(?:\.\d+)?)s\]\s*(.+)")


def parse_transcript_file(path: Path, time_offset: float = 0.0) -> list[tuple[float, str]]:
    """
    Legge un file di trascrizione con righe "[timestamp_float] testo".
    Applica time_offset a tutti i timestamp.
    """
    segments = []
    unmatched = 0
    for line in path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line:
            continue
        m = TRANSCRIPT_LINE.match(line)
        if m:
            t = float(m.group(1)) + time_offset
            segments.append((t, m.group(2).strip()))
        else:
            unmatched += 1
    if unmatched:
        print(f"  ⚠️  {path.name}: {unmatched} righe non riconosciute")
    return segments


def load_and_merge_transcripts(txt_paths: list[Path]) -> str:
    """
    Carica uno o più file di trascrizione, ordinati alfabeticamente.
    Per ogni file successivo al primo, aggiunge come offset il timestamp
    massimo del file precedente, rendendo la sequenza continua.
    Restituisce il testo plain concatenato (senza timestamp).
    """
    txt_paths = sorted(txt_paths, key=lambda p: p.name)
    all_segments: list[tuple[float, str]] = []
    offset = 0.0

    for path in txt_paths:
        segments = parse_transcript_file(path, time_offset=offset)
        if not segments:
            print(f"  ⚠️  Nessun segmento in {path.name}, saltato.")
            continue
        max_ts = max(t for t, _ in segments)
        print(f"   · {path.name}: {len(segments)} segmenti, ts max = {max_ts:.1f}s (offset applicato: {offset:.1f}s)")
        offset = max_ts
        all_segments.extend(segments)

    if not all_segments:
        return ""

    all_segments.sort(key=lambda x: x[0])
    return " ".join(text for _, text in all_segments)


# ─────────────────────────────────────────────
#  Chiamata Ollama
# ─────────────────────────────────────────────

SYSTEM_FULL = (
    "Sei un assistente che aiuta studenti universitari a studiare. "
    "Ricevi il testo estratto da un PowerPoint e la trascrizione della spiegazione del professore. "
    "Produci un riassunto strutturato, chiaro e utile allo studio in italiano, "
    "che integri entrambe le fonti. Organizza il riassunto per argomenti o sezioni logiche "
    "(non necessariamente slide per slide). Evidenzia i concetti chiave con grassetto."
)

SYSTEM_SLIDES_ONLY = (
    "Sei un assistente che aiuta studenti universitari a studiare. "
    "Ricevi il testo estratto da un PowerPoint. "
    "Produci un riassunto strutturato e chiaro in italiano, organizzato per argomenti. "
    "Evidenzia i concetti chiave con grassetto."
)


def call_ollama(model: str, system: str, user_text: str, images_b64: list[str], verbose: bool) -> str:
    """
    Chiama ollama.chat con testo e immagini opzionali.
    Le immagini vengono passate nel campo 'images' del messaggio user,
    come richiesto dall'API ollama-python per i modelli vision.
    """
    if verbose:
        print(f"   Invio a Ollama · modello: {model} · immagini: {len(images_b64)}")

    messages = [
        {"role": "system", "content": system},
        {"role": "user", "content": user_text, **({"images": images_b64} if images_b64 else {})},
    ]

    response = ollama.chat(model=model, messages=messages)
    return response["message"]["content"].strip()


# ─────────────────────────────────────────────
#  Entry point
# ─────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Riassume un PowerPoint con la trascrizione del professore via Ollama.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument("pptx", help="File .pptx")
    parser.add_argument(
        "transcripts", nargs="*", metavar="TRANSCRIPT",
        help="Uno o più file .txt di trascrizione (non usare con --slides-only)"
    )
    parser.add_argument("--slides-only", action="store_true",
                        help="Riassume solo le slide senza trascrizione")
    parser.add_argument("-m", "--model", default="gemma4",
                        help="Modello Ollama (default: gemma4). Available: 'qwen3.6'; online models: 'minimax-m2.7:cloud'")
    parser.add_argument("-o", "--output", default=None,
                        help="File di output .md (default: <nome_pptx>_riassunto.md)")
    parser.add_argument("-ni", "--no-images", action="store_true",
                        help="Non invia le immagini all'LLM")
    parser.add_argument("-mi", "--max-images", type=int, default=20,
                        help="Numero massimo di immagini uniche da estrarre (default: 20)")
    parser.add_argument("-v", "--verbose", action="store_true",
                        help="Output dettagliato")
    args = parser.parse_args()

    # ── Validazione argomenti ──
    if args.slides_only and args.transcripts:
        print("⚠️  --slides-only attivo: i file di trascrizione forniti verranno ignorati.")
    if not args.slides_only and not args.transcripts:
        sys.exit("Errore: fornisci almeno un file di trascrizione, oppure usa --slides-only.")

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        sys.exit(f"Errore: file non trovato: {pptx_path}")

    output_path = Path(args.output) if args.output else pptx_path.with_name(
        pptx_path.stem + "_riassunto.md"
    )

    # ── Estrazione testo slide / pagine ──
    suffix = pptx_path.suffix.lower()
    print(f"📂 Carico il file: {pptx_path.name}")
    if suffix == ".pdf":
        slide_text = extract_slide_text_pdf(pptx_path)
        extract_images_fn = extract_slide_images_pdf
    elif suffix in (".pptx", ".ppt"):
        slide_text = extract_slide_text(pptx_path)
        extract_images_fn = extract_slide_images
    else:
        sys.exit(f"Errore: formato non supportato '{suffix}'. Usa .pptx o .pdf")
    if args.verbose:
        print(f"   Testo estratto: {len(slide_text)} caratteri")

    # ── Estrazione immagini ──
    images_b64: list[str] = []
    if not args.no_images:
        print("🖼️  Estraggo immagini...")
        imgs = extract_images_fn(pptx_path, max_images=args.max_images)
        images_b64 = [img["b64"] for img in imgs]
        print(f"   → {len(images_b64)} immagini uniche estratte")
    else:
        print("🖼️  Immagini saltate (--no-images)")

    # ── Verifica Ollama ──
    try:
        ollama.list()
    except Exception as e:
        sys.exit(
            f"Errore: impossibile contattare Ollama ({e})\n"
            "Assicurati che sia in esecuzione con: ollama serve"
        )

    # ── Costruzione prompt e chiamata ──
    if args.slides_only:
        print(f"\n📋 Modalità: solo slide · modello: {args.model}")
        user_text = (
            f"## Testo delle slide\n\n{slide_text}\n\n"
            "Produci un riassunto strutturato per argomenti."
        )
        system = SYSTEM_SLIDES_ONLY
        meta = "solo slide"

    else:
        txt_paths = [Path(p) for p in args.transcripts]
        missing = [p for p in txt_paths if not p.exists()]
        if missing:
            sys.exit(f"Errore: file non trovati: {', '.join(str(p) for p in missing)}")

        print(f"\n📝 Carico {len(txt_paths)} file di trascrizione...")
        transcript = load_and_merge_transcripts(txt_paths)
        if not transcript:
            sys.exit("Errore: nessun testo estratto dai file di trascrizione.")
        print(f"   → Trascrizione totale: {len(transcript)} caratteri")

        print(f"\n📋 Modalità: slide + trascrizione · modello: {args.model}")
        user_text = (
            f"## Testo delle slide\n\n{slide_text}\n\n"
            f"## Trascrizione del professore\n\n{transcript}\n\n"
            "Produci un riassunto integrato e ben strutturato della lezione."
        )
        system = SYSTEM_FULL
        src_files = ", ".join(f"`{Path(p).name}`" for p in args.transcripts)
        meta = f"trascrizioni: {src_files}"

    print("✍️  Genero il riassunto...")
    try:
        summary = call_ollama(args.model, system, user_text, images_b64, args.verbose)
    except Exception as e:
        sys.exit(f"Errore Ollama: {e}")

    # ── Salva output ──
    output_md = (
        f"# Riassunto Lezione: {pptx_path.stem}\n\n"
        f"*Modello: `{args.model}` · {meta}*\n\n---\n\n"
        f"{summary}\n"
    )
    output_path.write_text(output_md, encoding="utf-8")
    print(f"\n✅ Salvato in: {output_path.resolve()}")


if __name__ == "__main__":
    main()